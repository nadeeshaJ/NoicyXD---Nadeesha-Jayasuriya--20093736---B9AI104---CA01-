"""Generate Phase 9 presentation slides, live demo script, and Q&A prep document.

Run:
    python scripts/generate_phase9_presentation.py

Outputs:
    reports/final/Presentation_Slides.pptx
    reports/final/Live_Demo_Script.docx
    reports/final/Presentation_Speaker_Notes.docx
"""

from __future__ import annotations

import json
import sys
from datetime import date
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PInches
from pptx.util import Pt as PPts

PROJECT_ROOT = Path(__file__).resolve().parents[1]
FINAL_DIR = PROJECT_ROOT / "reports" / "final"
FIG = PROJECT_ROOT / "reports" / "figures"


def load_ablation_summary() -> dict | None:
    path = FINAL_DIR / "ca1_ablation_summary.json"
    if path.exists():
        with path.open(encoding="utf-8") as f:
            return json.load(f)
    return None


def urban_metrics_for_slides() -> list[list[str]]:
    rows = []
    for model_name, label in (
        ("custom_cnn", "M1 Custom CNN"),
        ("resnet50", "M2 ResNet50"),
        ("mobilenetv2", "M3 MobileNetV2"),
    ):
        path = PROJECT_ROOT / "experiments" / "urbansound8k" / model_name / "test_metrics.json"
        with path.open(encoding="utf-8") as f:
            m = json.load(f)
        macro = m["classification_report"]["macro avg"]
        rows.append([
            label,
            f"{m['metrics']['accuracy']*100:.1f}%",
            f"{macro['recall']:.3f}",
            f"{m['metrics']['macro_f1']:.3f}",
        ])
    return rows

ACCENT = RGBColor(0x1E, 0x40, 0xAF)  # blue-800
DARK = RGBColor(0x1F, 0x29, 0x37)


def load_json(path: Path) -> dict:
    with path.open(encoding="utf-8") as f:
        return json.load(f)


def load_data() -> dict:
    return {
        "step3": load_json(PROJECT_ROOT / "reports" / "step3" / "step3_summary.json"),
        "step4": load_json(PROJECT_ROOT / "reports" / "step4" / "step4_summary.json"),
        "step5": load_json(PROJECT_ROOT / "reports" / "step5" / "step5_verification.json"),
        "step6": load_json(PROJECT_ROOT / "reports" / "step6" / "step6_summary.json"),
    }


def set_title_style(shape, size: int = 32, color=DARK):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            run.font.size = PPts(size)
            run.font.bold = True
            run.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    set_title_style(slide.shapes.title, 36)
    for p in slide.placeholders[1].text_frame.paragraphs:
        for run in p.runs:
            run.font.size = PPts(18)
            run.font.color.rgb = ACCENT


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    image_path: Path | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title, 28)

    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = PPts(18)
        for run in p.runs:
            run.font.color.rgb = DARK

    if image_path and image_path.exists():
        left = PInches(6.8)
        top = PInches(1.6)
        width = PInches(5.8)
        slide.shapes.add_picture(str(image_path), left, top, width=width)


def add_two_column_slide(
    prs: Presentation,
    title: str,
    left_bullets: list[str],
    right_bullets: list[str],
    left_heading: str = "",
    right_heading: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title, 28)

    def add_box(left, heading, items):
        tx = slide.shapes.add_textbox(left, PInches(1.5), PInches(5.8), PInches(5.0))
        tf = tx.text_frame
        tf.word_wrap = True
        if heading:
            hp = tf.paragraphs[0]
            hp.text = heading
            hp.font.bold = True
            hp.font.size = PPts(20)
            hp.font.color.rgb = ACCENT
        for i, item in enumerate(items):
            p = tf.add_paragraph() if heading or i > 0 else tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = PPts(16)
            p.space_after = PPts(6)
            for run in p.runs:
                run.font.color.rgb = DARK

    add_box(PInches(0.6), left_heading, left_bullets)
    add_box(PInches(6.8), right_heading, right_bullets)


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    image_path: Path | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    set_title_style(slide.shapes.title, 28)

    cols = len(headers)
    table_shape = slide.shapes.add_table(
        len(rows) + 1, cols, PInches(0.6), PInches(1.6), PInches(6.0 if image_path else 12.0), PInches(3.5)
    )
    table = table_shape.table
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = PPts(14)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = PPts(13)

    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), PInches(7.0), PInches(1.6), width=PInches(5.5))


def build_presentation(data: dict, ablation: dict | None = None) -> Presentation:
    s3, s4, s5, s6 = data["step3"], data["step4"], data["step5"], data["step6"]
    ablation = ablation or {}
    head_ablation = ablation.get("model3_head_ablation", [])
    best_urban_err = next(
        r for r in s6["urban_error_analysis"] if r["run_name"] == s6["best_urban_run"]
    )
    top_pair = best_urban_err["top_confused_pairs"][0]

    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    add_title_slide(
        prs,
        "Environmental Sound Classification\nUsing Mel-Spectrogram Images",
        "B9AI104 Deep Learning — Continuous Assessment\n"
        "UrbanSound8K + ESC-50 Cross-Domain Study\n"
        "Nadeesha Jayasuriya | Dr. Shahram Azizi Sazi | "
        f"{date.today().strftime('%B %Y')}",
    )

    add_bullet_slide(
        prs,
        "CA1 Brief Introduction (2–3 Sentences)",
        [
            "Task: multi-class environmental sound classification on Mel-spectrogram RGB images.",
            "Application: urban noise monitoring and animal sound recognition.",
            "Data: 8,732 UrbanSound8K clips (10 classes) + 200 ESC-50 animal clips (10 classes).",
            "Urban classes include: siren, jackhammer, dog_bark, street_music, etc.",
        ],
    )

    add_bullet_slide(
        prs,
        "Research Questions",
        [
            "Can CNNs classify environmental sounds from Mel-spectrogram images?",
            "How do Custom CNN, ResNet50, and MobileNetV2 compare?",
            "Does the same pipeline generalise from urban to animal sounds?",
            "Which transfer-learning strategy works on small ESC-50 data?",
            "Which model is best for real-time deployment?",
        ],
    )

    add_bullet_slide(
        prs,
        "Why Convert Sound → Image?",
        [
            "CNNs excel at spatial patterns — Mel-specs show frequency over time",
            "Transfer learning: reuse ImageNet-pretrained ResNet / MobileNet",
            "Fixed 224×224 input aligns with standard vision architectures",
            "Alternative MFCC vectors need different models (1D CNN / RNN)",
            "Mel scale matches human hearing — good for environmental sounds",
        ],
        FIG / "step6" / "mel_vs_mfcc_comparison.png",
    )

    add_two_column_slide(
        prs,
        "Datasets",
        [
            "8,732 clips, 10 urban classes",
            "Official fold-10 test split",
            "Train / val / test: 7105 / 790 / 837",
            "Classes: siren, jackhammer, dog_bark, etc.",
        ],
        [
            "200 clips, 10 animal classes (HF subset)",
            "Stratified 70 / 15 / 15 split",
            "Train / val / test: 140 / 30 / 30",
            "Classes: dog, cow, frog, crow, etc.",
        ],
        left_heading="Stage 1: UrbanSound8K",
        right_heading="Stage 2: ESC-50 Animals",
    )

    add_bullet_slide(
        prs,
        "Preprocessing Pipeline",
        [
            "Load WAV → resample 22,050 Hz mono",
            "Pad / trim to 4 seconds",
            "STFT → Mel filterbank (128 bins) → dB scale",
            "Min-max normalise → 224×224 RGB PNG",
            "Same config for training, evaluation, and Streamlit app",
        ],
        FIG / "step2" / "preprocessing_pipeline_diagram.png",
    )

    add_bullet_slide(
        prs,
        "CA1 Three-Model Structure",
        [
            "Model 1: Custom CNN — conventional baseline from scratch (conv, pool, dense, dropout)",
            "Model 2: ResNet50 — traditional ImageNet transfer with 10-class head",
            "Model 3: MobileNetV2 — customised transfer head + fine-tune (best model)",
            "Model 1 is NOT the CA1 customised model — Model 3 holds that role",
            "Compared by accuracy, macro recall, macro F1, and training time",
            "Model 3 head-size ablation: tested 64, 128, 256 hidden nodes",
        ],
    )

    add_bullet_slide(
        prs,
        "Three Model Architectures",
        [
            "Model 1 — conventional CNN baseline from scratch (6.7M params)",
            "Model 2 — ResNet50 ImageNet transfer, residual blocks (23.5M params)",
            "Model 3 — MobileNetV2 efficient depthwise separable convs (2.2M params)",
            "Transfer models: freeze backbone → train head → fine-tune",
            "Loss: cross-entropy | Optimiser: Adam | LR: 1e-4",
        ],
        FIG / "step3" / "architecture_transfer_learning.png",
    )

    add_table_slide(
        prs,
        "UrbanSound8K — Accuracy & Recall (Fold 10)",
        ["Model", "Accuracy", "Macro recall", "Macro F1"],
        urban_metrics_for_slides(),
        FIG / "step3" / "model_comparison_urbansound8k.png",
    )

    add_bullet_slide(
        prs,
        "Training Curves & Early Stopping",
        [
            "Best checkpoint selected by minimum validation loss",
            f"Early stopping patience = 5 epochs (no val improvement)",
            "Watch val loss — not training loss alone (avoid overfitting)",
            "MobileNetV2 best model: 82.7% test accuracy, 0.831 macro F1",
        ],
        FIG / "step3" / "urbansound8k" / "mobilenetv2" / "training_history.png",
    )

    if head_ablation:
        ab_rows = []
        for row in head_ablation:
            ab_rows.append([
                str(row.get("hidden_units", "")),
                f"{row.get('test_accuracy', 0)*100:.1f}%",
                f"{row.get('test_macro_recall', 0):.3f}",
            ])
        add_table_slide(
            prs,
            "Model 3 Head-Size Ablation (8 epochs)",
            ["Hidden units", "Accuracy", "Macro recall"],
            ab_rows,
            FIG / "final" / "ablation_model3_head_sizes.png",
        )

    esc_rows = []
    for run in s4["esc50_runs"]:
        label = run["run_name"].replace("_", " ")
        esc_rows.append([
            label[:28],
            f"{run['test_metrics']['accuracy']*100:.1f}%",
            f"{run['test_metrics']['macro_f1']:.3f}",
        ])
    add_bullet_slide(
        prs,
        "ESC-50 Transfer Learning",
        [
            f"Best approach: {s4['best_esc50_run']} (F1 = {s4['best_esc50_macro_f1']:.3f})",
            "Urban transfer: F1 = 0.516 — domain shift hurts urban features",
            "From scratch CNN: F1 = 0.513 — insufficient data (140 train clips)",
            "ImageNet MobileNetV2: F1 = 0.607 — best generalisation",
            "Honest finding: ImageNet beat urban pretraining on animals",
        ],
    )
    add_table_slide(prs, "ESC-50 Comparison Table", ["Approach", "Acc.", "Macro F1"], esc_rows)

    cross_rows = []
    for row in s4["cross_domain"]:
        cross_rows.append([
            row["approach"][:24],
            f"{row['macro_f1']:.3f}",
            f"{row['f1_drop']*100:.1f}%",
        ])
    add_table_slide(
        prs,
        "Cross-Domain Performance Drop",
        ["Approach", "Macro F1", "F1 Drop"],
        cross_rows,
        FIG / "step4" / "cross_domain_comparison.png",
    )

    add_bullet_slide(
        prs,
        "Error Analysis Example",
        [
            f"Top urban confusion: {top_pair['true_class']} → {top_pair['predicted_class']} ({top_pair['count']}×)",
            "Siren vs children playing — outdoor ambient overlap",
            "Drilling vs jackhammer — similar mechanical rhythms",
            "Dog bark vs street music — transient bursts in urban scenes",
            "Errors are acoustically explainable, not random",
        ],
        FIG / "step6" / "case_studies" / "case_study_01.png",
    )

    urban_check = s5["checks"][0]
    add_bullet_slide(
        prs,
        "Streamlit Live Demo",
        [
            "Two modes: Urban Sound / Animal Sound",
            "Upload WAV → waveform + Mel-spec + 224×224 preview",
            "Top-3 predictions with confidence bars",
            f"Urban: dog_bark at {urban_check['confidence']:.0%} (demo clip)",
            "Run: python -m streamlit run app/streamlit_app.py",
        ],
        FIG / "step5" / "app_demo_urban.png",
    )

    add_bullet_slide(
        prs,
        "Q&A Concepts (Know These)",
        [
            "Weight sharing: same CNN filter slides across the image → translation invariance",
            "Backprop: cross-entropy loss → gradients → Adam weight updates each epoch",
            "ML vs DL: DL has hidden conv/dense layers; logistic regression has none",
            "Probabilities: Streamlit shows confidence % — not just class names",
            "If Model 3 ablation < full model: full run used longer fine-tuning + early stopping",
        ],
    )

    add_bullet_slide(
        prs,
        "Conclusion & Limitations",
        [
            "Best urban model: MobileNetV2 — 82.7% acc, 0.831 macro F1",
            "Selected for deployment: highest F1 + smallest checkpoint (8.8 MB)",
            "Custom CNN is faster at inference (~0.9 ms) but lower accuracy — trade-off explained",
            "Pipeline generalises; domain shift remains the main challenge",
            "Limitations: small ESC-50 test set, fixed 4 s window, CPU Stage-1 training",
            "Future: more data, audio augmentation, domain-adaptive fine-tuning",
        ],
    )

    add_title_slide(
        prs,
        "Thank You",
        "Questions?\n\nDemo ready: http://localhost:8501",
    )

    return prs


def add_heading(doc: Document, text: str, level: int = 1):
    return doc.add_heading(text, level=level)


def add_para(doc: Document, text: str, bold: bool = False):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.font.size = Pt(11)


def add_bullets(doc: Document, items: list[str]):
    for item in items:
        p = doc.add_paragraph(item, style="List Bullet")
        for run in p.runs:
            run.font.size = Pt(11)


def add_numbered(doc: Document, items: list[str]):
    for item in items:
        p = doc.add_paragraph(item, style="List Number")
        for run in p.runs:
            run.font.size = Pt(11)


def build_demo_script(data: dict) -> Document:
    s5 = data["step5"]
    s6 = data["step6"]
    urban_audio = Path(s5["checks"][0]["sample_audio"])
    animal_audio = Path(s5["checks"][1]["sample_audio"])
    challenge_audio = PROJECT_ROOT / "data/raw/urbansound8k/audio/fold10/165166-8-0-2.wav"

    doc = Document()
    title = doc.add_heading("Live Demo Script — Presentation", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    add_para(doc, "Phase 9: Environmental Sound Classification Streamlit Demo")
    doc.add_paragraph()

    add_heading(doc, "Before You Start (2 minutes)", level=1)
    add_bullets(
        doc,
        [
            "Open terminal in project folder: noicy_XD",
            "Run: python -m streamlit run app/streamlit_app.py",
            "Confirm browser opens at http://localhost:8501",
            "Have these 3 WAV files ready in File Explorer:",
            f"  1. Urban easy: {urban_audio}",
            f"  2. Animal easy: {animal_audio}",
            f"  3. Urban challenge: {challenge_audio}",
        ],
    )

    add_heading(doc, "Demo Flow (~5 minutes)", level=1)

    add_para(doc, "Step 1 — Introduce the app (30 sec)", bold=True)
    add_para(
        doc,
        "Say: 'This Streamlit app runs the same preprocessing and models used in our experiments. "
        "Users upload a WAV file and get Mel-spectrogram visualisation plus top-3 class predictions.'",
    )

    add_para(doc, "Step 2 — Urban mode, easy correct case (90 sec)", bold=True)
    add_numbered(
        doc,
        [
            "Sidebar: select Urban Sound Mode",
            "Upload file: 100795-3-0-0.wav (dog bark)",
            "Point to waveform and Mel-spectrogram panels",
            "Highlight 224×224 RGB model input",
            f"Show prediction: dog_bark at ~{s5['checks'][0]['confidence']:.0%} confidence",
            "Say: 'MobileNetV2 achieved 82.7% accuracy on fold-10 urban test set.'",
        ],
    )

    add_para(doc, "Step 3 — Animal mode (90 sec)", bold=True)
    add_numbered(
        doc,
        [
            "Switch sidebar to Animal Sound Mode (model reloads)",
            f"Upload: {animal_audio.name}",
            "Show lower confidence is expected on 30-test-clip ESC-50 subset",
            f"Top prediction should be dog (~{s5['checks'][1]['confidence']:.0%})",
            "Say: 'ImageNet transfer outperformed urban pretraining on animal sounds.'",
        ],
    )

    add_para(doc, "Step 4 — Challenging misclassification (90 sec)", bold=True)
    add_numbered(
        doc,
        [
            "Switch back to Urban Sound Mode",
            "Upload: 165166-8-0-2.wav (true class: siren)",
            "If predicted as children_playing, explain from Step 6 case study:",
            "  'Distant sirens mixed with outdoor noise resemble children playing in Mel images.'",
            "Say: 'This shows our error analysis — confusions are acoustically similar classes.'",
        ],
    )

    add_para(doc, "Step 5 — Wrap up (30 sec)", bold=True)
    add_para(
        doc,
        "Say: 'We recommend MobileNetV2 for deployment — best urban F1 and smallest model size. "
        "Custom CNN has lower latency but lower accuracy. Main limitation is cross-domain shift "
        "when moving from urban to animal sounds.'",
    )

    add_heading(doc, "Backup If Upload Fails", level=1)
    add_bullets(
        doc,
        [
            "Show static screenshots in reports/figures/step5/app_demo_urban.png",
            "Walk through Final_Assignment_Report.docx Section 6 (Deployment)",
            "Show case study figure: reports/figures/step6/case_studies/case_study_01.png",
        ],
    )

    add_heading(doc, "Anticipated Questions & Answers", level=1)
    qa = [
        (
            "Why Mel-spectrogram instead of raw audio?",
            "CNNs need grid-structured input. Mel-specs preserve harmonic and transient patterns "
            "as 2D images and allow ImageNet transfer learning.",
        ),
        (
            "Why 224×224?",
            "Standard input for ResNet/MobileNet pretrained on ImageNet; keeps deployment simple.",
        ),
        (
            "Why transfer learning for ESC-50?",
            "Only 140 training clips — training from scratch gives F1 ≈ 0.51; ImageNet transfer reaches 0.61.",
        ),
        (
            "Which model would you deploy and why?",
            "MobileNetV2: highest urban macro F1 (0.831) and smallest checkpoint (8.8 MB). "
            "Custom CNN is faster (~0.9 ms) but less accurate — MobileNetV2 is the better "
            "accuracy-efficiency trade-off.",
        ),
        (
            "What would you improve with more time?",
            "More animal data, stronger augmentation, domain adaptation, and app UI polish with demo clips built in.",
        ),
        (
            "Why did urban transfer fail on animals?",
            "Urban spectrogram features (sirens, jackhammers) do not transfer to farm/animal vocalisations — domain shift.",
        ),
    ]
    for q, a in qa:
        add_para(doc, f"Q: {q}", bold=True)
        add_para(doc, f"A: {a}")
        doc.add_paragraph()

    return doc


def build_speaker_notes(data: dict) -> Document:
    doc = Document()
    doc.add_heading("Presentation Speaker Notes", level=0)
    add_para(doc, "12 slides — approx. 10–12 minute talk + 5 minute demo")

    notes = [
        ("Slide 1 — Title", "Introduce yourself, module, and project title. State the two-stage design."),
        ("Slide 2 — Research Questions", "Frame the talk: we are not just building a classifier — we compare models and domains."),
        ("Slide 3 — Sound to Image", "Emphasise Mel vs MFCC figure if shown. This is the core design choice."),
        ("Slide 4 — Datasets", "Stress fold-10 official test for urban; small ESC-50 for cross-domain only."),
        ("Slide 5 — Pipeline", "Mention config.yaml ties preprocess, train, and app together."),
        ("Slide 6 — CA1 Model Structure", "Stress: M1 = conventional baseline; M2/M3 = transfer; M3 = customised."),
        ("Slide 7 — Architectures", "Three models = conventional baseline + strong transfer + efficient customised transfer."),
        ("Slide 8 — Urban Results", "MobileNetV2 wins. Custom CNN proves transfer learning value."),
        ("Slide 9 — Training Curves", "Best checkpoint by val loss; early stopping patience = 5."),
        ("Slide 10 — ESC-50", "Be honest: ImageNet beat urban transfer — interesting negative result."),
        ("Slide 11 — Cross-Domain", "22–31% F1 drop quantifies domain shift."),
        ("Slide 12 — Error Analysis", "Show case study — examiners like explainable failures."),
        ("Slide 13 — Demo", "Transition to live Streamlit or screenshots if live demo risky."),
        ("Slide 14 — Q&A Concepts", "Know weight sharing, backprop, model.summary parameter counts."),
        ("Slide 15 — Conclusion", "MobileNetV2 = best F1 + smallest size; Custom CNN faster but weaker."),
    ]
    for slide, note in notes:
        add_para(doc, slide, bold=True)
        add_para(doc, note)
        doc.add_paragraph()

    return doc


def main() -> None:
    FINAL_DIR.mkdir(parents=True, exist_ok=True)
    data = load_data()
    ablation = load_ablation_summary()

    prs = build_presentation(data, ablation)
    slides_path = FINAL_DIR / "Presentation_Slides.pptx"
    prs.save(slides_path)

    demo_script = build_demo_script(data)
    demo_path = FINAL_DIR / "Live_Demo_Script.docx"
    demo_script.save(demo_path)

    speaker_notes = build_speaker_notes(data)
    notes_path = FINAL_DIR / "Presentation_Speaker_Notes.docx"
    speaker_notes.save(notes_path)

    manifest = {
        "phase": 9,
        "title": "Presentation & Demo Prep",
        "generated": date.today().isoformat(),
        "slide_count": len(prs.slides),
        "outputs": {
            "slides": str(slides_path),
            "demo_script": str(demo_path),
            "speaker_notes": str(notes_path),
        },
    }
    with (FINAL_DIR / "phase9_manifest.json").open("w", encoding="utf-8") as f:
        json.dump(manifest, f, indent=2)

    print("Phase 9 deliverables generated:")
    print(f"  Slides ({len(prs.slides)} slides): {slides_path}")
    print(f"  Demo script:          {demo_path}")
    print(f"  Speaker notes:        {notes_path}")


if __name__ == "__main__":
    main()
